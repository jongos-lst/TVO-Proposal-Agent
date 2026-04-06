"""Generate a professional TVO proposal PowerPoint deck with Getac corporate branding.

10-slide structure:
  1. Cover
  2. Executive Summary (KPIs + value proposition)
  3. Customer Challenges
  4. Recommended Solution (per product)
  5. TCO Analysis — data table + chart (per product)
  6. ROI & Savings
  7. Risk & Reliability
  8. Competitive Differentiation
  9. Conclusion & Next Steps
 10. Thank You

Brand: Getac orange #D74B00 accent, charcoal headers, white slide backgrounds.
All numeric data sourced from TVOCalculation — no hallucinated values.
"""

import os
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

from app.models.proposal import Proposal
from app.services.chart_generator import generate_all_charts

# ── Getac Corporate Brand Colors ─────────────────────────────────
GETAC_ORANGE = RGBColor(0xD7, 0x4B, 0x00)
GETAC_DARK_ORANGE = RGBColor(0xB8, 0x3D, 0x00)
CHARCOAL = RGBColor(0x1B, 0x1B, 0x2F)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xF8)
BORDER_GRAY = RGBColor(0xE2, 0xE4, 0xE8)
GETAC_BLUE = RGBColor(0x00, 0x66, 0xB3)
COMPETITOR_RED = RGBColor(0xB8, 0x35, 0x2A)
SAVINGS_GREEN = RGBColor(0x1A, 0x87, 0x54)
MUTED_TEXT = RGBColor(0x6B, 0x70, 0x80)

SLIDE_W = Inches(13.333)  # Widescreen 16:9
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────

def _add_accent_bar(slide, x, y, width, height, color=GETAC_ORANGE):
    """Add a thin accent bar / decorative line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text(slide, x, y, w, h, text, size=16, color=DARK_TEXT, bold=False,
              align=PP_ALIGN.LEFT):
    """Helper to add a text box."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def _add_header_bar(slide, title: str):
    """Add a consistent header bar to content slides (charcoal bg + orange accent)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                   SLIDE_W, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CHARCOAL
    shape.line.fill.background()

    _add_accent_bar(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04),
                    GETAC_ORANGE)

    _add_text(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
              title, size=28, color=WHITE, bold=True)


def _add_footer(slide, page_num: int):
    """Add a subtle footer with page number and branding."""
    _add_accent_bar(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.02),
                    BORDER_GRAY)
    _add_text(slide, Inches(0.8), Inches(7.15), Inches(4), Inches(0.3),
              "GETAC  |  Total Value of Ownership Proposal", size=8,
              color=MUTED_TEXT)
    _add_text(slide, Inches(11), Inches(7.15), Inches(1.5), Inches(0.3),
              f"{page_num}", size=8, color=MUTED_TEXT, align=PP_ALIGN.RIGHT)


def _add_kpi_card(slide, x, y, w, h, value: str, label: str,
                  value_color=GETAC_ORANGE):
    """Add a metric/KPI card with a large value and a label beneath."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = BORDER_GRAY
    card.line.width = Pt(0.5)

    _add_text(slide, Inches(x + 0.15), Inches(y + 0.12), Inches(w - 0.3),
              Inches(h * 0.5), value, size=22, color=value_color, bold=True,
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(x + 0.15), Inches(y + h * 0.55), Inches(w - 0.3),
              Inches(h * 0.35), label, size=9, color=DARK_TEXT,
              align=PP_ALIGN.CENTER)


def _format_cell(cell, text: str, size: int = 10, color=DARK_TEXT,
                 bold: bool = False, fill=None, align=PP_ALIGN.LEFT):
    """Format a table cell with consistent styling."""
    cell.text = text
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def _add_productivity_breakdown_table(slide, factors, deployment_years: int,
                                      left, top, width, row_h=0.3):
    """Add a compact table showing per-factor productivity value."""
    from app.models.tvo import ProductivityFactor
    rows = len(factors) + 2  # header + factors + total
    cols = 3
    table_h = Inches(rows * row_h)
    shape = slide.shapes.add_table(rows, cols, left, top, width, table_h)
    table = shape.table

    # Column widths
    table.columns[0].width = Inches(5.0)
    table.columns[1].width = Inches(3.65)
    table.columns[2].width = Inches(3.65)

    # Header
    header_fill = RGBColor(0x2C, 0x2C, 0x2C)
    for ci, hdr in enumerate(["Productivity Factor", "Annual Value", f"{deployment_years}-Year Total"]):
        _format_cell(table.cell(0, ci), hdr, size=9, bold=True,
                     color=RGBColor(0xFF, 0xFF, 0xFF), fill=header_fill,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # Factor rows
    total_annual = 0.0
    total_total = 0.0
    for ri, factor in enumerate(factors, start=1):
        _format_cell(table.cell(ri, 0), factor.name, size=9)
        _format_cell(table.cell(ri, 1), f"${factor.annual_value:,.0f}", size=9,
                     align=PP_ALIGN.CENTER)
        _format_cell(table.cell(ri, 2), f"${factor.total_value:,.0f}", size=9,
                     align=PP_ALIGN.CENTER, bold=True, color=SAVINGS_GREEN)
        total_annual += factor.annual_value
        total_total += factor.total_value
        # Alternate row shading
        if ri % 2 == 0:
            for ci in range(cols):
                table.cell(ri, ci).fill.solid()
                table.cell(ri, ci).fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    # Total row
    total_ri = len(factors) + 1
    _format_cell(table.cell(total_ri, 0), "TOTAL PRODUCTIVITY VALUE", size=9,
                 bold=True, fill=RGBColor(0xE8, 0xF5, 0xE9))
    _format_cell(table.cell(total_ri, 1), f"${total_annual:,.0f}", size=9,
                 bold=True, align=PP_ALIGN.CENTER,
                 fill=RGBColor(0xE8, 0xF5, 0xE9))
    _format_cell(table.cell(total_ri, 2), f"${total_total:,.0f}", size=10,
                 bold=True, align=PP_ALIGN.CENTER, color=SAVINGS_GREEN,
                 fill=RGBColor(0xE8, 0xF5, 0xE9))


def _serialize_products(products):
    """Serialize a list of GetacProduct objects to dicts."""
    if not products:
        return []
    return [p.model_dump() if hasattr(p, "model_dump") else p for p in products]


def _serialize_tvo_results(tvo_results):
    """Serialize a dict of product_id -> TVOCalculation to dicts."""
    if not tvo_results:
        return {}
    return {k: (v.model_dump() if hasattr(v, "model_dump") else v)
            for k, v in tvo_results.items()}


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 1: Cover
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _add_title_slide(prs: Presentation, proposal: Proposal):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = CHARCOAL

    # Orange accent bar at top
    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12),
                    GETAC_ORANGE)

    # GETAC wordmark (top-left)
    _add_text(slide, Inches(1.0), Inches(0.5), Inches(4), Inches(0.6),
              "GETAC", size=18, color=GETAC_ORANGE, bold=True)

    # Main title
    _add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
              "Total Value of Ownership", size=44, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER)

    # Subtitle
    _add_text(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.6),
              "Proposal & TCO Analysis", size=24, color=GETAC_ORANGE,
              align=PP_ALIGN.CENTER)

    # Orange divider
    _add_accent_bar(slide, Inches(5.5), Inches(3.8), Inches(2.333), Inches(0.03),
                    GETAC_ORANGE)

    # Customer name
    customer = (proposal.persona.customer_name or "Valued Customer"
                if proposal.persona else "Valued Customer")
    _add_text(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(0.6),
              f"Prepared for: {customer}", size=20, color=WHITE,
              align=PP_ALIGN.CENTER)

    # Industry
    if proposal.persona and proposal.persona.industry:
        _add_text(slide, Inches(1.5), Inches(4.7), Inches(10), Inches(0.35),
                  f"Industry: {proposal.persona.industry}", size=14,
                  color=RGBColor(0x99, 0x9E, 0xAA), align=PP_ALIGN.CENTER)

    # Product list
    if proposal.selected_products:
        product_names = ", ".join(p.name for p in proposal.selected_products)
        _add_text(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(0.35),
                  f"Products: {product_names}", size=13,
                  color=RGBColor(0x99, 0x9E, 0xAA), align=PP_ALIGN.CENTER)

    # Date
    _add_text(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.35),
              datetime.now().strftime("%B %d, %Y"), size=14,
              color=RGBColor(0x99, 0x9E, 0xAA), align=PP_ALIGN.CENTER)

    # Bottom bar
    _add_accent_bar(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.5),
                    RGBColor(0x12, 0x12, 0x20))
    _add_text(slide, Inches(1), Inches(7.05), Inches(11), Inches(0.4),
              "GETAC  \u2022  Rugged Computing Solutions", size=11,
              color=RGBColor(0x99, 0x9E, 0xAA), align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 2: Executive Summary
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _add_executive_summary(prs: Presentation, proposal: Proposal, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "Executive Summary")
    _add_footer(slide, page_num)

    tvo_results = proposal.tvo_calculations
    if not tvo_results:
        _add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                  "TVO calculation data not yet available.",
                  size=16, color=MUTED_TEXT, align=PP_ALIGN.CENTER)
        return

    # Aggregate metrics across all products
    total_savings = sum(t.tco_savings for t in tvo_results.values())
    total_comp_tco = sum(t.competitor_total_tco for t in tvo_results.values())
    savings_pct = (total_savings / total_comp_tco * 100) if total_comp_tco > 0 else 0
    total_productivity = sum(t.productivity_savings_total for t in tvo_results.values())
    avg_risk_reduction = (sum(t.risk_reduction_percent for t in tvo_results.values())
                          / len(tvo_results))
    first_tvo = next(iter(tvo_results.values()))
    roi_label = (f"Month {first_tvo.roi_payback_months:.0f}"
                 if first_tvo.roi_payback_months < first_tvo.deployment_years * 12
                 else "N/A")

    # Section label
    _add_text(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.4),
              "KEY PERFORMANCE INDICATORS", size=11, color=GETAC_ORANGE,
              bold=True)
    _add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(2), Inches(0.03),
                    GETAC_ORANGE)

    # 4 KPI cards in a row
    card_w = 2.7
    card_h = 1.3
    gap = 0.4
    start_x = (13.333 - (card_w * 4 + gap * 3)) / 2
    kpis = [
        (f"${total_savings:,.0f}", f"TCO Savings ({savings_pct:.1f}%)", SAVINGS_GREEN),
        (roi_label, "ROI Break-even", GETAC_ORANGE),
        (f"{avg_risk_reduction:.0f}%", "Risk Reduction", GETAC_BLUE),
        (f"${total_productivity:,.0f}", "Productivity Savings", SAVINGS_GREEN),
    ]
    for i, (value, label, color) in enumerate(kpis):
        x = start_x + i * (card_w + gap)
        _add_kpi_card(slide, x, 2.1, card_w, card_h, value, label, color)

    # Value proposition
    if proposal.value_proposition:
        _add_text(slide, Inches(0.8), Inches(3.8), Inches(5), Inches(0.4),
                  "VALUE PROPOSITION", size=11, color=GETAC_ORANGE, bold=True)
        _add_accent_bar(slide, Inches(0.8), Inches(4.15), Inches(2), Inches(0.03),
                        GETAC_ORANGE)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8), Inches(4.4),
                                      Inches(11.7), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = BORDER_GRAY
        card.line.width = Pt(0.5)

        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(4.55),
                                         Inches(11.3), Inches(1.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = proposal.value_proposition
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.line_spacing = Pt(20)

    # Fleet info footnote
    _add_text(slide, Inches(0.8), Inches(6.4), Inches(11), Inches(0.3),
              f"Based on {first_tvo.fleet_size} units over "
              f"{first_tvo.deployment_years}-year deployment  |  "
              f"{len(tvo_results)} product(s) analyzed",
              size=9, color=MUTED_TEXT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 3: Customer Situation & Pain Points
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _add_pain_points_slide(prs: Presentation, proposal: Proposal, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "Customer Situation & Challenges")
    _add_footer(slide, page_num)

    persona = proposal.persona
    if not persona:
        return

    # Left column: Pain Points
    _add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
              "Key Pain Points", size=18, color=CHARCOAL, bold=True)
    _add_accent_bar(slide, Inches(0.8), Inches(1.85), Inches(1.5), Inches(0.03),
                    GETAC_ORANGE)

    y = 2.1
    if persona.pain_points:
        for point in persona.pain_points:
            # Orange left-border accent bar
            _add_accent_bar(slide, Inches(0.8), Inches(y), Inches(0.06),
                            Inches(0.5), GETAC_ORANGE)
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(0.86), Inches(y),
                                          Inches(5.44), Inches(0.5))
            card.fill.solid()
            card.fill.fore_color.rgb = LIGHT_BG
            card.line.fill.background()
            _add_text(slide, Inches(1.05), Inches(y + 0.07), Inches(5.1),
                      Inches(0.36), point, size=12, color=DARK_TEXT)
            y += 0.6

    # Right column: Current Environment
    _add_text(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.5),
              "Current Environment", size=18, color=CHARCOAL, bold=True)
    _add_accent_bar(slide, Inches(7), Inches(1.85), Inches(1.5), Inches(0.03),
                    GETAC_ORANGE)

    y_right = 2.1
    if persona.current_devices:
        _add_text(slide, Inches(7.2), Inches(y_right), Inches(5), Inches(0.4),
                  "Current Devices:", size=13, color=DARK_TEXT, bold=True)
        y_right += 0.4
        for dev in persona.current_devices:
            _add_text(slide, Inches(7.4), Inches(y_right), Inches(5),
                      Inches(0.35), f"\u2022  {dev}", size=12, color=DARK_TEXT)
            y_right += 0.35

    if persona.use_scenarios:
        y_right += 0.3
        _add_text(slide, Inches(7.2), Inches(y_right), Inches(5), Inches(0.4),
                  "Use Scenarios:", size=13, color=DARK_TEXT, bold=True)
        y_right += 0.4
        for scenario in persona.use_scenarios:
            _add_text(slide, Inches(7.4), Inches(y_right), Inches(5),
                      Inches(0.35), f"\u2022  {scenario}", size=12,
                      color=DARK_TEXT)
            y_right += 0.35

    if persona.fleet_size:
        y_right += 0.3
        _add_text(slide, Inches(7.2), Inches(y_right), Inches(5), Inches(0.35),
                  f"Fleet Size: {persona.fleet_size} units", size=13,
                  color=GETAC_ORANGE, bold=True)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Per-product: Recommended Solution slide
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _add_solution_slide(prs: Presentation, product, product_index: int,
                        total_products: int, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = f"Recommended Solution \u2014 {product.name}"
    if total_products > 1:
        title = f"Solution {product_index + 1} of {total_products} \u2014 {product.name}"
    _add_header_bar(slide, title)
    _add_footer(slide, page_num)

    # Product name + category badge
    _add_text(slide, Inches(0.8), Inches(1.4), Inches(7), Inches(0.6),
              product.name, size=28, color=CHARCOAL, bold=True)
    _add_text(slide, Inches(0.8), Inches(1.95), Inches(4), Inches(0.35),
              product.category.upper(), size=11, color=GETAC_ORANGE, bold=True)

    _add_accent_bar(slide, Inches(0.8), Inches(2.4), Inches(2), Inches(0.03),
                    GETAC_ORANGE)

    # Specs table
    specs = [
        ("Display", product.display_size),
        ("Processor", product.processor),
        ("Rugged Rating", product.rugged_rating),
        ("Battery", product.battery_life),
        ("Warranty", product.warranty_standard),
        ("Price", f"${product.base_price:,.0f}"),
    ]

    rows = len(specs)
    table_shape = slide.shapes.add_table(rows, 2, Inches(0.8), Inches(2.6),
                                         Inches(5.5), Inches(0.38 * rows))
    table = table_shape.table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(3.7)

    for i, (label, value) in enumerate(specs):
        cell_label = table.cell(i, 0)
        _format_cell(cell_label, label, size=11, color=CHARCOAL, bold=True,
                     fill=LIGHT_BG)

        cell_val = table.cell(i, 1)
        row_fill = WHITE if i % 2 == 0 else RGBColor(0xFD, 0xFD, 0xFD)
        _format_cell(cell_val, value, size=11, color=DARK_TEXT, fill=row_fill)

    # Key features (right side)
    _add_text(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.5),
              "Key Features", size=18, color=CHARCOAL, bold=True)
    _add_accent_bar(slide, Inches(7), Inches(1.85), Inches(1.5), Inches(0.03),
                    GETAC_ORANGE)

    y_feat = 2.1
    if product.key_features:
        for feat in product.key_features[:6]:
            _add_text(slide, Inches(7.2), Inches(y_feat), Inches(5),
                      Inches(0.38), f"\u2713  {feat}", size=12,
                      color=DARK_TEXT)
            y_feat += 0.4

    if product.target_industries:
        y_feat += 0.2
        _add_text(slide, Inches(7), Inches(y_feat), Inches(5.5), Inches(0.35),
                  "Target Industries", size=13, color=CHARCOAL, bold=True)
        y_feat += 0.35
        _add_text(slide, Inches(7.2), Inches(y_feat), Inches(5), Inches(0.35),
                  " \u2022 ".join(product.target_industries), size=11,
                  color=DARK_TEXT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Per-product: TCO Analysis slide (DATA TABLE + chart)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _add_tvo_slide(prs: Presentation, product, tvo, charts: dict[str, str],
                   page_num: int, total_products: int, product_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = f"TCO Analysis \u2014 {product.name}"
    if total_products > 1:
        title = f"TCO Analysis ({product_index + 1}/{total_products}) \u2014 {product.name}"
    _add_header_bar(slide, title)
    _add_footer(slide, page_num)

    if not tvo:
        return

    # ── LEFT: TCO Data Table ──
    _add_text(slide, Inches(0.6), Inches(1.25), Inches(4), Inches(0.35),
              "COST BREAKDOWN", size=10, color=GETAC_ORANGE, bold=True)

    num_items = len(tvo.tco_line_items)
    table_rows = num_items + 2  # header + items + total row
    table_shape = slide.shapes.add_table(
        table_rows, 4, Inches(0.6), Inches(1.6), Inches(7.0), Inches(0.36 * table_rows)
    )
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(1.6)
    table.columns[2].width = Inches(1.6)
    table.columns[3].width = Inches(1.6)

    # Header row
    headers = ["Category", "Getac", "Competitor", "Savings"]
    header_colors = [CHARCOAL, GETAC_BLUE, COMPETITOR_RED, SAVINGS_GREEN]
    for col_idx, (hdr, hdr_color) in enumerate(zip(headers, header_colors)):
        cell = table.cell(0, col_idx)
        _format_cell(cell, hdr, size=10, color=WHITE, bold=True,
                     fill=CHARCOAL, align=PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT)

    # Data rows
    for i, item in enumerate(tvo.tco_line_items):
        row_idx = i + 1
        row_fill = WHITE if i % 2 == 0 else LIGHT_BG

        _format_cell(table.cell(row_idx, 0), item.label, size=10,
                     color=DARK_TEXT, fill=row_fill)
        _format_cell(table.cell(row_idx, 1), f"${item.getac_value:,.0f}",
                     size=10, color=DARK_TEXT, fill=row_fill,
                     align=PP_ALIGN.RIGHT)
        _format_cell(table.cell(row_idx, 2), f"${item.competitor_value:,.0f}",
                     size=10, color=DARK_TEXT, fill=row_fill,
                     align=PP_ALIGN.RIGHT)

        # Savings: green if positive (savings), red if negative (Getac more expensive)
        diff = item.difference
        diff_color = SAVINGS_GREEN if diff >= 0 else COMPETITOR_RED
        diff_text = f"${diff:,.0f}" if diff >= 0 else f"-${abs(diff):,.0f}"
        _format_cell(table.cell(row_idx, 3), diff_text, size=10,
                     color=diff_color, bold=True, fill=row_fill,
                     align=PP_ALIGN.RIGHT)

    # Total row
    total_row = table_rows - 1
    total_fill = RGBColor(0xE8, 0xEB, 0xF0)
    _format_cell(table.cell(total_row, 0), "TOTAL", size=10, color=CHARCOAL,
                 bold=True, fill=total_fill)
    _format_cell(table.cell(total_row, 1), f"${tvo.getac_total_tco:,.0f}",
                 size=10, color=GETAC_BLUE, bold=True, fill=total_fill,
                 align=PP_ALIGN.RIGHT)
    _format_cell(table.cell(total_row, 2), f"${tvo.competitor_total_tco:,.0f}",
                 size=10, color=COMPETITOR_RED, bold=True, fill=total_fill,
                 align=PP_ALIGN.RIGHT)
    savings_color = SAVINGS_GREEN if tvo.tco_savings >= 0 else COMPETITOR_RED
    savings_text = f"${tvo.tco_savings:,.0f}" if tvo.tco_savings >= 0 else f"-${abs(tvo.tco_savings):,.0f}"
    _format_cell(table.cell(total_row, 3), savings_text, size=10,
                 color=savings_color, bold=True, fill=total_fill,
                 align=PP_ALIGN.RIGHT)

    # ── RIGHT: Total TCO bar chart ──
    # Chart figsize is 5x4 (1.25:1 aspect ratio) — preserve it
    total_chart = charts.get("total_tco")
    if total_chart and os.path.exists(total_chart):
        chart_w = Inches(4.8)
        chart_h = Inches(3.84)  # 4.8 / 1.25 = 3.84 preserves 5:4 ratio
        slide.shapes.add_picture(total_chart, Inches(8.0), Inches(1.3),
                                 chart_w, chart_h)

    # ── BOTTOM: Summary stat cards ──
    # y_stats must clear both the table bottom and chart bottom
    table_bottom = 1.6 + 0.36 * table_rows
    chart_bottom = (1.3 + 3.84) if (total_chart and os.path.exists(total_chart)) else 0
    y_stats = max(table_bottom, chart_bottom) + 0.2
    # Cap to leave room for 1.0" cards + footnote before footer at 7.1
    y_stats = min(y_stats, 5.8)
    stats = [
        (f"${tvo.getac_total_tco:,.0f}", "Getac Total TCO", GETAC_BLUE),
        (f"${tvo.competitor_total_tco:,.0f}", "Competitor TCO", COMPETITOR_RED),
        (f"${tvo.tco_savings:,.0f}", "Total Savings", SAVINGS_GREEN),
        (f"{tvo.tco_savings_percent:.1f}%", "TCO Reduction", GETAC_ORANGE),
    ]

    card_w = 2.5
    gap = 0.3
    start_x = (13.333 - (card_w * 4 + gap * 3)) / 2
    for i, (value, label, color) in enumerate(stats):
        x = start_x + i * (card_w + gap)
        _add_kpi_card(slide, x, y_stats, card_w, 1.0, value, label, color)

    # Assumptions footnote
    num_factors = sum(1 for f in tvo.productivity_breakdown if f.applies) if tvo.productivity_breakdown else 0
    factor_note = f", {num_factors}-factor productivity model" if num_factors > 0 else ""
    _add_text(slide, Inches(0.6), Inches(6.7), Inches(12), Inches(0.3),
              f"Assumptions: {tvo.fleet_size} units, {tvo.deployment_years}-year "
              f"deployment{factor_note}",
              size=8, color=MUTED_TEXT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Combined: ROI & Savings
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _add_savings_slide(prs: Presentation, proposal: Proposal,
                       charts: dict[str, str], page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "ROI & Savings Analysis")
    _add_footer(slide, page_num)

    tvo_results = proposal.tvo_calculations
    if not tvo_results:
        return

    # Top row: ROI timeline (left) + Savings breakdown (right)
    # ROI chart figsize is 7x4.5 (1.556:1 aspect ratio) — preserve it
    roi_chart = charts.get("roi_timeline")
    if roi_chart and os.path.exists(roi_chart):
        roi_w = Inches(6.2)
        roi_h = Inches(3.98)  # 6.2 / 1.556 = 3.98 preserves 7:4.5 ratio
        slide.shapes.add_picture(roi_chart, Inches(0.3), Inches(1.3),
                                 roi_w, roi_h)

    # Savings breakdown figsize is 5x4 (1.25:1 aspect ratio) — preserve it
    savings_chart = charts.get("savings_breakdown")
    if savings_chart and os.path.exists(savings_chart):
        sav_w = Inches(5.0)
        sav_h = Inches(4.0)  # 5.0 / 1.25 = 4.0 preserves 5:4 ratio
        slide.shapes.add_picture(savings_chart, Inches(7.3), Inches(1.3),
                                 sav_w, sav_h)

    # Combined metrics across all products
    total_productivity_savings = sum(t.productivity_savings_total
                                     for t in tvo_results.values())
    total_tco_savings = sum(t.tco_savings for t in tvo_results.values())
    total_value = sum(t.total_value_advantage for t in tvo_results.values())

    first_tvo = next(iter(tvo_results.values()))
    roi_label = (f"Month {first_tvo.roi_payback_months:.0f}"
                 if first_tvo.roi_payback_months < first_tvo.deployment_years * 12
                 else "N/A")

    y_cards = 5.6
    metrics = [
        (f"${total_productivity_savings:,.0f}", "Productivity Savings",
         SAVINGS_GREEN),
        (roi_label, "ROI Break-even", GETAC_ORANGE),
        (f"${total_tco_savings:,.0f}", "Combined TCO Savings", GETAC_BLUE),
        (f"${total_value:,.0f}", "Total Value Advantage", SAVINGS_GREEN),
    ]

    card_w = 2.7
    gap = 0.35
    start_x = (13.333 - (card_w * 4 + gap * 3)) / 2
    for i, (value, label, color) in enumerate(metrics):
        x = start_x + i * (card_w + gap)
        _add_kpi_card(slide, x, y_cards, card_w, 1.1, value, label, color)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Risk & Reliability Analysis
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _add_risk_slide(prs: Presentation, proposal: Proposal,
                    charts: dict[str, str], page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "Risk & Reliability Analysis")
    _add_footer(slide, page_num)

    tvo_results = proposal.tvo_calculations
    if not tvo_results:
        return

    # Check if productivity breakdown data exists (determines chart dimensions)
    first_tvo = next(iter(tvo_results.values()), None)
    has_breakdown = bool(
        first_tvo and first_tvo.productivity_breakdown
        and any(f.applies and f.total_value > 0 for f in first_tvo.productivity_breakdown)
    )

    # Productivity / downtime chart — full-width when 3-panel, left half when 2-panel
    # Chart figsize: 12x3.5 (3.429:1) with breakdown, 8x3.5 (2.286:1) without
    prod_chart = charts.get("productivity")
    if prod_chart and os.path.exists(prod_chart):
        if has_breakdown:
            prod_w = Inches(12.5)
            prod_h = Inches(3.05)  # Shorter to leave room for cards + table
            slide.shapes.add_picture(prod_chart, Inches(0.3), Inches(1.3),
                                     prod_w, prod_h)
        else:
            prod_w = Inches(7.0)
            prod_h = Inches(3.06)  # 7.0 / 2.286 = 3.06
            slide.shapes.add_picture(prod_chart, Inches(0.3), Inches(1.3),
                                     prod_w, prod_h)

    # Cost waterfall chart (right) — only when no breakdown (to save space)
    if not has_breakdown:
        waterfall_chart = charts.get("cost_waterfall")
        if waterfall_chart and os.path.exists(waterfall_chart):
            wf_w = Inches(5.3)
            wf_h = Inches(2.98)  # 5.3 / 1.778 = 2.98
            slide.shapes.add_picture(waterfall_chart, Inches(7.7), Inches(1.3),
                                     wf_w, wf_h)

    # Aggregate stats
    total_getac_failures = sum(t.getac_expected_failures
                               for t in tvo_results.values())
    total_comp_failures = sum(t.competitor_expected_failures
                              for t in tvo_results.values())
    avg_risk_reduction = (sum(t.risk_reduction_percent
                              for t in tvo_results.values()) / len(tvo_results))
    total_downtime_saved = sum(
        (t.competitor_annual_downtime_hours - t.getac_annual_downtime_hours)
        * t.deployment_years
        for t in tvo_results.values()
    )
    total_prod_savings = sum(t.productivity_savings_total
                             for t in tvo_results.values())

    # Layout order: chart (1.3–4.35) → KPI cards → table → footer (7.1)
    # KPI stat cards
    card_h = 0.9 if has_breakdown else 1.1
    y_stats = 4.50 if has_breakdown else 4.8
    stats = [
        (f"{avg_risk_reduction:.0f}%", "Risk Reduction", SAVINGS_GREEN),
        (f"{total_getac_failures:.0f}", "Getac Failures (Expected)",
         GETAC_BLUE),
        (f"{total_comp_failures:.0f}", "Competitor Failures (Expected)",
         COMPETITOR_RED),
        (f"{total_downtime_saved:,.0f}h", "Downtime Hours Saved",
         SAVINGS_GREEN),
    ]

    card_w = 2.7
    gap = 0.35
    start_x = (13.333 - (card_w * 4 + gap * 3)) / 2
    for i, (value, label, color) in enumerate(stats):
        x = start_x + i * (card_w + gap)
        _add_kpi_card(slide, x, y_stats, card_w, card_h, value, label, color)

    # Productivity breakdown table below KPI cards (compact 0.2" rows)
    if has_breakdown and first_tvo:
        applicable = [f for f in first_tvo.productivity_breakdown
                      if f.applies and f.total_value > 0]
        if applicable:
            table_top = y_stats + card_h + 0.10
            _add_productivity_breakdown_table(
                slide, applicable, first_tvo.deployment_years,
                left=Inches(0.5), top=Inches(table_top),
                width=Inches(12.3), row_h=0.2,
            )


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Competitive Differentiation
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _add_competitive_slide(prs: Presentation, proposal: Proposal,
                           page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "Competitive Differentiation")
    _add_footer(slide, page_num)

    # Collect all advantages across all products
    all_advantages = []
    for product in proposal.selected_products:
        prod_advantages = proposal.competitive_advantages.get(product.id, [])
        competitor_name = proposal.competitor_product_names.get(product.id,
                                                                "Competitor")
        for adv in prod_advantages:
            label = (f"{product.name} vs {competitor_name}"
                     if len(proposal.selected_products) > 1 else "")
            all_advantages.append((adv, label))

    # Deduplicate
    seen = set()
    unique_advantages = []
    for adv, label in all_advantages:
        if adv not in seen:
            seen.add(adv)
            unique_advantages.append((adv, label))

    col_items = len(unique_advantages) // 2 + len(unique_advantages) % 2

    y = 1.5
    for i, (adv, label) in enumerate(unique_advantages):
        x = 0.8 if i < col_items else 7.0
        row_y = y + (i % col_items) * 0.75

        # Orange left-border accent
        _add_accent_bar(slide, Inches(x), Inches(row_y), Inches(0.06),
                        Inches(0.58), GETAC_ORANGE)

        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(x + 0.06), Inches(row_y),
                                      Inches(5.44), Inches(0.58))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.fill.background()

        _add_text(slide, Inches(x + 0.25), Inches(row_y + 0.1), Inches(5),
                  Inches(0.38), f"\u2713  {adv}", size=12, color=DARK_TEXT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Conclusion & Next Steps
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _add_conclusion_slide(prs: Presentation, proposal: Proposal,
                          page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "Conclusion & Next Steps")
    _add_footer(slide, page_num)

    # Value proposition
    if proposal.value_proposition:
        _add_text(slide, Inches(0.8), Inches(1.35), Inches(3), Inches(0.3),
                  "VALUE PROPOSITION", size=10, color=GETAC_ORANGE, bold=True)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8), Inches(1.7),
                                      Inches(11.7), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = GETAC_ORANGE
        card.line.width = Pt(1.5)

        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8),
                                         Inches(11.3), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = proposal.value_proposition
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.line_spacing = Pt(21)

    # Key numbers summary (left column)
    tvo_results = proposal.tvo_calculations
    if tvo_results:
        total_savings = sum(t.tco_savings for t in tvo_results.values())
        total_productivity = sum(t.productivity_savings_total
                                 for t in tvo_results.values())
        avg_risk = (sum(t.risk_reduction_percent for t in tvo_results.values())
                    / len(tvo_results))
        total_comp_tco = sum(t.competitor_total_tco
                             for t in tvo_results.values())
        savings_pct = ((total_savings / total_comp_tco * 100)
                       if total_comp_tco > 0 else 0)

        y_nums = 3.5
        prod_label = (f"KEY NUMBERS ({len(tvo_results)} PRODUCT"
                      f"{'S' if len(tvo_results) > 1 else ''})")
        _add_text(slide, Inches(0.8), Inches(y_nums), Inches(5), Inches(0.4),
                  prod_label, size=10, color=GETAC_ORANGE, bold=True)
        _add_accent_bar(slide, Inches(0.8), Inches(y_nums + 0.3),
                        Inches(1.5), Inches(0.03), GETAC_ORANGE)

        numbers = [
            (f"${total_savings:,.0f}", "Combined TCO Savings"),
            (f"{savings_pct:.1f}%", "Cost Reduction"),
            (f"${total_productivity:,.0f}", "Productivity Value"),
            (f"{avg_risk:.0f}%", "Avg Risk Reduction"),
        ]
        for i, (val, label) in enumerate(numbers):
            y_n = y_nums + 0.5 + i * 0.45
            _add_text(slide, Inches(1.0), Inches(y_n), Inches(2.5),
                      Inches(0.35), val, size=16, color=SAVINGS_GREEN,
                      bold=True)
            _add_text(slide, Inches(3.5), Inches(y_n), Inches(4), Inches(0.35),
                      label, size=13, color=DARK_TEXT)

    # Next steps (right column)
    next_steps = [
        "Schedule a hands-on product demo",
        "Conduct a pilot program with a subset of the fleet",
        "Finalize pricing and warranty terms",
        "Plan deployment timeline and training",
    ]

    _add_text(slide, Inches(7), Inches(3.5), Inches(5.5), Inches(0.4),
              "RECOMMENDED NEXT STEPS", size=10, color=GETAC_ORANGE,
              bold=True)
    _add_accent_bar(slide, Inches(7), Inches(3.8), Inches(1.5), Inches(0.03),
                    GETAC_ORANGE)

    y_step = 4.0
    for i, step in enumerate(next_steps, 1):
        # Orange numbered circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.15),
                                        Inches(y_step + 0.02), Inches(0.35),
                                        Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = GETAC_ORANGE
        circle.line.fill.background()
        _add_text(slide, Inches(7.15), Inches(y_step + 0.04), Inches(0.35),
                  Inches(0.3), str(i), size=12, color=WHITE, bold=True,
                  align=PP_ALIGN.CENTER)

        _add_text(slide, Inches(7.7), Inches(y_step + 0.02), Inches(5),
                  Inches(0.35), step, size=13, color=DARK_TEXT)
        y_step += 0.5

    # Closing statement
    _add_accent_bar(slide, Inches(0.8), Inches(6.2), Inches(11.7),
                    Inches(0.02), BORDER_GRAY)
    _add_text(slide, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.4),
              "Thank you for considering Getac. We look forward to partnering with you.",
              size=14, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 10: Thank You
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _add_thank_you_slide(prs: Presentation, proposal: Proposal):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = CHARCOAL

    # Orange accent bar at top
    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12),
                    GETAC_ORANGE)

    # "Thank You" large text
    _add_text(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(1.0),
              "Thank You", size=48, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER)

    # Orange divider
    _add_accent_bar(slide, Inches(5.5), Inches(3.5), Inches(2.333),
                    Inches(0.03), GETAC_ORANGE)

    # Customer reference
    customer = (proposal.persona.customer_name or "Valued Customer"
                if proposal.persona else "Valued Customer")
    _add_text(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.5),
              f"Prepared for {customer}", size=16,
              color=RGBColor(0x99, 0x9E, 0xAA), align=PP_ALIGN.CENTER)

    # Date
    _add_text(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.5),
              datetime.now().strftime("%B %d, %Y"), size=14,
              color=RGBColor(0x99, 0x9E, 0xAA), align=PP_ALIGN.CENTER)

    # GETAC wordmark
    _add_text(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.6),
              "GETAC", size=24, color=GETAC_ORANGE, bold=True,
              align=PP_ALIGN.CENTER)

    # Bottom bar
    _add_accent_bar(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.5),
                    RGBColor(0x12, 0x12, 0x20))
    _add_text(slide, Inches(1), Inches(7.05), Inches(11), Inches(0.4),
              "GETAC  \u2022  Rugged Computing Solutions", size=11,
              color=RGBColor(0x99, 0x9E, 0xAA), align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Main generator (multi-product)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def generate_proposal_pptx(proposal: Proposal,
                           output_dir: str | None = None) -> BytesIO:
    """Generate a complete TVO proposal PowerPoint deck (multi-product).

    Returns a BytesIO buffer containing the .pptx file.
    """
    prs = Presentation()
    prs.slide_width = Emu(int(13.333 * 914400))  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    chart_dir = output_dir or os.path.join(os.path.dirname(__file__), "..",
                                           "output", "charts")

    # Generate charts per product
    product_charts: dict[str, dict[str, str]] = {}
    first_charts: dict[str, str] = {}
    for i, product in enumerate(proposal.selected_products):
        tvo = proposal.tvo_calculations.get(product.id)
        if tvo:
            prod_chart_dir = (os.path.join(chart_dir, product.id)
                              if len(proposal.selected_products) > 1
                              else chart_dir)
            charts = generate_all_charts(tvo, prod_chart_dir)
            product_charts[product.id] = charts
            if i == 0:
                first_charts = charts

    # Build slides
    page = 1

    # Slide 1: Cover
    _add_title_slide(prs, proposal)
    page += 1

    # Slide 2: Executive Summary
    _add_executive_summary(prs, proposal, page)
    page += 1

    # Slide 3: Customer Situation
    _add_pain_points_slide(prs, proposal, page)
    page += 1

    # Per-product: Solution + TCO slides
    total_products = len(proposal.selected_products)
    for i, product in enumerate(proposal.selected_products):
        _add_solution_slide(prs, product, i, total_products, page)
        page += 1

        tvo = proposal.tvo_calculations.get(product.id)
        charts = product_charts.get(product.id, {})
        if tvo:
            _add_tvo_slide(prs, product, tvo, charts, page, total_products, i)
            page += 1

    # ROI & Savings (uses first product's charts for visuals)
    _add_savings_slide(prs, proposal, first_charts, page)
    page += 1

    # Risk & Reliability
    _add_risk_slide(prs, proposal, first_charts, page)
    page += 1

    # Competitive differentiation
    _add_competitive_slide(prs, proposal, page)
    page += 1

    # Conclusion
    _add_conclusion_slide(prs, proposal, page)
    page += 1

    # Thank You
    _add_thank_you_slide(prs, proposal)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
