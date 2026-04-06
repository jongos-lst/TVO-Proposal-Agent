"""Integration tests for PowerPoint generation (Getac corporate theme)."""

import sys
import os

import pytest
from pptx import Presentation

from app.services.pptx_generator import generate_proposal_pptx

# Reuse sample proposal builder
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
from scripts.generate_sample_proposal import build_sample_proposal


@pytest.fixture
def sample_proposal():
    return build_sample_proposal()


def _extract_all_text(prs: Presentation) -> str:
    """Extract all text from all slides into a single string."""
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    all_text.append(p.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            all_text.append(p.text)
    return " ".join(all_text)


class TestPptxGeneration:
    def test_generates_valid_pptx(self, sample_proposal, tmp_path):
        buf = generate_proposal_pptx(sample_proposal, output_dir=str(tmp_path))
        prs = Presentation(buf)
        # Single product: 10 slides (cover, exec summary, challenges,
        # solution, TCO, savings, risk, competitive, conclusion, thank you)
        assert len(prs.slides) == 10

    def test_slide_titles_present(self, sample_proposal, tmp_path):
        buf = generate_proposal_pptx(sample_proposal, output_dir=str(tmp_path))
        prs = Presentation(buf)
        full_text = _extract_all_text(prs)
        assert "Total Value of Ownership" in full_text
        assert "Executive Summary" in full_text
        assert "Customer Situation" in full_text
        assert "Risk & Reliability" in full_text
        assert "Competitive Differentiation" in full_text
        assert "Conclusion" in full_text
        assert "Thank You" in full_text

    def test_customer_name_in_deck(self, sample_proposal, tmp_path):
        buf = generate_proposal_pptx(sample_proposal, output_dir=str(tmp_path))
        prs = Presentation(buf)
        full_text = _extract_all_text(prs)
        assert "Pacific Northwest Utilities" in full_text

    def test_charts_embedded(self, sample_proposal, tmp_path):
        """Verify that picture shapes (charts) are embedded."""
        buf = generate_proposal_pptx(sample_proposal, output_dir=str(tmp_path))
        prs = Presentation(buf)
        picture_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    picture_count += 1
        # At least 4 charts (TCO, total_tco, savings, ROI, productivity; waterfall
        # is omitted when productivity breakdown table is present)
        assert picture_count >= 4

    def test_pptx_returns_bytesio(self, sample_proposal, tmp_path):
        from io import BytesIO
        buf = generate_proposal_pptx(sample_proposal, output_dir=str(tmp_path))
        assert isinstance(buf, BytesIO)
        assert buf.getbuffer().nbytes > 0

    def test_tco_table_present(self, sample_proposal, tmp_path):
        """Verify TCO data table has line item labels."""
        buf = generate_proposal_pptx(sample_proposal, output_dir=str(tmp_path))
        prs = Presentation(buf)
        full_text = _extract_all_text(prs)
        assert "Hardware Acquisition" in full_text
        assert "Extended Warranty" in full_text
        assert "Repair & Replacement" in full_text
        assert "TOTAL" in full_text

    def test_executive_summary_kpis(self, sample_proposal, tmp_path):
        """Verify executive summary slide has KPI values."""
        buf = generate_proposal_pptx(sample_proposal, output_dir=str(tmp_path))
        prs = Presentation(buf)
        full_text = _extract_all_text(prs)
        # Should contain savings percentage and risk reduction
        assert "TCO Savings" in full_text
        assert "Risk Reduction" in full_text
        assert "ROI Break-even" in full_text
        assert "Productivity Savings" in full_text

    def test_no_unintentional_shape_overlaps(self, sample_proposal, tmp_path):
        """Verify no shapes overlap unintentionally on any slide.

        Allowed overlaps: text boxes fully contained inside a larger shape
        (e.g. header text inside header bar, KPI label inside card).
        """
        buf = generate_proposal_pptx(sample_proposal, output_dir=str(tmp_path))
        prs = Presentation(buf)

        for slide_idx, slide in enumerate(prs.slides):
            shapes = []
            for s in slide.shapes:
                l = s.left
                t = s.top
                r = s.left + s.width
                b = s.top + s.height
                shapes.append((s.name, l, t, r, b))

            for i in range(len(shapes)):
                for j in range(i + 1, len(shapes)):
                    n1, l1, t1, r1, b1 = shapes[i]
                    n2, l2, t2, r2, b2 = shapes[j]
                    # Check bounding box intersection
                    if l1 < r2 and r1 > l2 and t1 < b2 and b1 > t2:
                        # Allow if one is fully contained in the other
                        contained = (
                            (l2 >= l1 and r2 <= r1 and t2 >= t1 and b2 <= b1) or
                            (l1 >= l2 and r1 <= r2 and t1 >= t2 and b1 <= b2)
                        )
                        if not contained:
                            overlap_w = min(r1, r2) - max(l1, l2)
                            overlap_h = min(b1, b2) - max(t1, t2)
                            area_emu2 = overlap_w * overlap_h
                            # Allow small overlaps: adjacent text stacking, KPI
                            # card value/label pairs (< 0.5 sq inches)
                            small_threshold = 0.5 * (914400 ** 2)
                            if area_emu2 >= small_threshold:
                                assert False, (
                                    f"Slide {slide_idx + 1}: {n1} overlaps {n2}"
                                )
